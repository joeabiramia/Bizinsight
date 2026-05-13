"""Executive Report Generation Service.

Generates enterprise-grade reports in:
  - PDF (via reportlab)
  - PowerPoint (via python-pptx)
  - Excel (via openpyxl)
"""
from __future__ import annotations

import io
from datetime import datetime
from typing import Any

import pandas as pd

from app.dataframe_utils import safe_number

# ── Color palette ──────────────────────────────────────────────────────────────
_BRAND_BLUE = (0.094, 0.333, 0.718)   # #184fb7
_BRAND_DARK = (0.094, 0.102, 0.18)    # #181a2e
_ACCENT_GREEN = (0.047, 0.698, 0.510) # #0CB282
_ACCENT_RED = (0.906, 0.224, 0.224)   # #E73939
_LIGHT_GRAY = (0.94, 0.94, 0.96)
_WHITE = (1, 1, 1)
_TEXT_DARK = (0.1, 0.1, 0.15)


# ─────────────────────────────────────────────────────────────────────────────
# PDF REPORT
# ─────────────────────────────────────────────────────────────────────────────

def generate_pdf_report(
    df: pd.DataFrame,
    analysis: dict,
    insights: list[dict],
    file_id: str,
    company_name: str = "BizInsight AI",
) -> bytes:
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import cm, mm
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
            HRFlowable, PageBreak,
        )
        from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
    except ImportError:
        raise ImportError("reportlab not installed. Run: pip install reportlab")

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        rightMargin=2*cm, leftMargin=2*cm,
        topMargin=2*cm, bottomMargin=2*cm
    )

    brand_blue = colors.Color(*_BRAND_BLUE)
    brand_dark = colors.Color(*_BRAND_DARK)
    accent_green = colors.Color(*_ACCENT_GREEN)
    light_gray = colors.Color(*_LIGHT_GRAY)

    styles = getSampleStyleSheet()
    style_h1 = ParagraphStyle("h1", parent=styles["Heading1"], fontSize=22, textColor=brand_blue, spaceAfter=6)
    style_h2 = ParagraphStyle("h2", parent=styles["Heading2"], fontSize=14, textColor=brand_dark, spaceAfter=4, spaceBefore=10)
    style_h3 = ParagraphStyle("h3", parent=styles["Heading3"], fontSize=11, textColor=brand_dark, spaceAfter=3, spaceBefore=6)
    style_body = ParagraphStyle("body", parent=styles["Normal"], fontSize=9, textColor=colors.Color(*_TEXT_DARK), spaceAfter=4, leading=14)
    style_caption = ParagraphStyle("caption", parent=styles["Normal"], fontSize=8, textColor=colors.gray, italic=True)
    style_center = ParagraphStyle("center", parent=style_body, alignment=TA_CENTER)
    style_kpi = ParagraphStyle("kpi", parent=styles["Normal"], fontSize=11, textColor=brand_blue, fontName="Helvetica-Bold")

    generated_at = datetime.utcnow().strftime("%B %d, %Y")
    industry = analysis.get("industry", "Business")
    shape = analysis.get("shape", {})
    rows = shape.get("rows", len(df))
    cols_count = shape.get("columns", len(df.columns))

    story = []

    # ── Cover Page ────────────────────────────────────────────────────────────
    story.append(Spacer(1, 3*cm))
    story.append(Paragraph(company_name, ParagraphStyle("cover_title", parent=style_h1, fontSize=28, alignment=TA_CENTER)))
    story.append(Spacer(1, 0.4*cm))
    story.append(Paragraph("AI Executive Business Report", ParagraphStyle("cover_sub", parent=style_h2, fontSize=18, alignment=TA_CENTER, textColor=colors.gray)))
    story.append(Spacer(1, 0.8*cm))
    story.append(HRFlowable(width="80%", thickness=2, color=brand_blue, hAlign="CENTER"))
    story.append(Spacer(1, 0.8*cm))
    story.append(Paragraph(f"Industry: {industry}", style_center))
    story.append(Paragraph(f"Dataset: {rows:,} records × {cols_count} columns", style_center))
    story.append(Paragraph(f"Generated: {generated_at}", style_center))
    story.append(PageBreak())

    # ── Executive Summary ─────────────────────────────────────────────────────
    story.append(Paragraph("1. Executive Summary", style_h1))
    story.append(HRFlowable(width="100%", thickness=1, color=brand_blue))
    story.append(Spacer(1, 0.3*cm))

    num_insights = len(insights)
    num_critical = sum(1 for i in insights if i.get("type") == "risk")
    num_opportunities = sum(1 for i in insights if i.get("type") == "opportunity")
    story.append(Paragraph(
        f"This report presents an AI-powered analysis of {rows:,} business records across {cols_count} dimensions. "
        f"The analysis identified <b>{num_insights} key business insights</b>, including "
        f"<b>{num_critical} risk factors</b> and <b>{num_opportunities} growth opportunities</b>. "
        f"Industry classification: <b>{industry}</b>.",
        style_body
    ))
    story.append(Spacer(1, 0.4*cm))

    # ── KPI Highlights ────────────────────────────────────────────────────────
    story.append(Paragraph("2. KPI Highlights", style_h1))
    story.append(HRFlowable(width="100%", thickness=1, color=brand_blue))
    story.append(Spacer(1, 0.3*cm))

    numeric_summary = analysis.get("numeric_summary", {})
    if numeric_summary:
        kpi_data = [["Metric", "Total", "Mean", "Median", "Std Dev", "Max"]]
        for col, stats in list(numeric_summary.items())[:8]:
            kpi_data.append([
                col[:20],
                f"{stats.get('total', 0):,.2f}",
                f"{stats.get('mean', 0):,.2f}",
                f"{stats.get('median', 0):,.2f}",
                f"{stats.get('std', 0):,.2f}",
                f"{stats.get('max', 0):,.2f}",
            ])

        kpi_table = Table(kpi_data, colWidths=[4*cm, 2.8*cm, 2.8*cm, 2.8*cm, 2.8*cm, 2.8*cm])
        kpi_table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), brand_blue),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 8),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.Color(0.96, 0.97, 1.0)]),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.Color(0.85, 0.85, 0.9)),
            ("ALIGN", (1, 0), (-1, -1), "RIGHT"),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("TOPPADDING", (0, 0), (-1, -1), 5),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
        ]))
        story.append(kpi_table)
    story.append(Spacer(1, 0.5*cm))

    # ── Key Findings & Insights ───────────────────────────────────────────────
    story.append(Paragraph("3. Key Findings", style_h1))
    story.append(HRFlowable(width="100%", thickness=1, color=brand_blue))
    story.append(Spacer(1, 0.3*cm))

    for i, insight in enumerate(insights[:8], 1):
        insight_type = insight.get("type", "info")
        color_map = {"risk": _ACCENT_RED, "opportunity": _ACCENT_GREEN, "performance": _BRAND_BLUE}
        marker_color = colors.Color(*color_map.get(insight_type, _BRAND_BLUE))

        story.append(Paragraph(f"<b>{i}. {insight.get('title', 'Finding')}</b>", style_h3))
        if insight.get("observation"):
            story.append(Paragraph(f"<b>Observation:</b> {insight['observation']}", style_body))
        if insight.get("interpretation"):
            story.append(Paragraph(f"<b>Interpretation:</b> {insight['interpretation']}", style_body))
        if insight.get("action"):
            story.append(Paragraph(f"<b>Recommended Action:</b> {insight['action']}", style_body))
        story.append(Spacer(1, 0.2*cm))

    # ── Risks ─────────────────────────────────────────────────────────────────
    risk_insights = [i for i in insights if i.get("type") == "risk"]
    if risk_insights:
        story.append(PageBreak())
        story.append(Paragraph("4. Risk Assessment", style_h1))
        story.append(HRFlowable(width="100%", thickness=1, color=brand_blue))
        story.append(Spacer(1, 0.3*cm))
        for i, risk in enumerate(risk_insights, 1):
            story.append(Paragraph(f"<b>Risk {i}: {risk.get('title', '')}</b>", style_h3))
            story.append(Paragraph(risk.get("observation", ""), style_body))
            story.append(Paragraph(f"<i>Mitigation: {risk.get('action', 'Monitor closely.')}</i>", style_caption))
            story.append(Spacer(1, 0.2*cm))

    # ── Opportunities ─────────────────────────────────────────────────────────
    opp_insights = [i for i in insights if i.get("type") == "opportunity"]
    if opp_insights:
        story.append(Paragraph("5. Growth Opportunities", style_h1))
        story.append(HRFlowable(width="100%", thickness=1, color=brand_blue))
        story.append(Spacer(1, 0.3*cm))
        for i, opp in enumerate(opp_insights, 1):
            story.append(Paragraph(f"<b>Opportunity {i}: {opp.get('title', '')}</b>", style_h3))
            story.append(Paragraph(opp.get("observation", ""), style_body))
            story.append(Paragraph(f"<i>Action: {opp.get('action', '')}</i>", style_caption))
            story.append(Spacer(1, 0.2*cm))

    # ── Recommendations ───────────────────────────────────────────────────────
    story.append(Paragraph("6. AI Recommendations", style_h1))
    story.append(HRFlowable(width="100%", thickness=1, color=brand_blue))
    story.append(Spacer(1, 0.3*cm))
    story.append(Paragraph(
        "Based on the AI analysis of your business data, the following strategic recommendations are prioritized:",
        style_body
    ))
    story.append(Spacer(1, 0.2*cm))

    all_actions = [i.get("action", "") for i in insights if i.get("action")]
    for idx, action in enumerate(all_actions[:6], 1):
        story.append(Paragraph(f"<b>{idx}.</b> {action}", style_body))

    # ── Footer ────────────────────────────────────────────────────────────────
    story.append(Spacer(1, 1*cm))
    story.append(HRFlowable(width="100%", thickness=0.5, color=colors.gray))
    story.append(Paragraph(f"Generated by BizInsight AI | {generated_at} | Confidential", style_caption))

    doc.build(story)
    return buf.getvalue()


# ─────────────────────────────────────────────────────────────────────────────
# POWERPOINT REPORT
# ─────────────────────────────────────────────────────────────────────────────

def generate_pptx_report(
    df: pd.DataFrame,
    analysis: dict,
    insights: list[dict],
    file_id: str,
    company_name: str = "BizInsight AI",
) -> bytes:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        raise ImportError("python-pptx not installed. Run: pip install python-pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blue = RGBColor(24, 79, 183)
    dark = RGBColor(24, 26, 46)
    green = RGBColor(12, 178, 130)
    red = RGBColor(231, 57, 57)
    white = RGBColor(255, 255, 255)
    light = RGBColor(240, 242, 255)

    generated_at = datetime.utcnow().strftime("%B %d, %Y")
    industry = analysis.get("industry", "Business")
    shape = analysis.get("shape", {})
    rows = shape.get("rows", len(df))

    blank_layout = prs.slide_layouts[6]

    def _add_rect(slide, left, top, width, height, fill_rgb):
        from pptx.util import Inches
        shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
        shape.line.fill.background()
        return shape

    def _add_text_box(slide, text, left, top, width, height, font_size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
        from pptx.util import Inches, Pt
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        return txBox

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _add_rect(slide, 0, 0, 13.33, 7.5, dark)
    _add_rect(slide, 0, 0, 13.33, 0.08, blue)
    _add_rect(slide, 0, 7.42, 13.33, 0.08, blue)
    _add_text_box(slide, company_name, 1, 2, 11, 1.2, 36, True, white, PP_ALIGN.CENTER)
    _add_text_box(slide, "AI Executive Business Report", 1, 3.2, 11, 0.8, 20, False, RGBColor(160, 175, 220), PP_ALIGN.CENTER)
    _add_text_box(slide, f"Industry: {industry}  |  {rows:,} Records  |  {generated_at}", 1, 4.2, 11, 0.6, 12, False, RGBColor(130, 145, 180), PP_ALIGN.CENTER)

    # ── Slide 2: Executive Summary ────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _add_rect(slide, 0, 0, 13.33, 1.1, blue)
    _add_text_box(slide, "Executive Summary", 0.4, 0.2, 10, 0.8, 24, True, white)
    _add_text_box(slide, f"Dataset: {rows:,} records  |  Industry: {industry}  |  {generated_at}", 0.4, 6.9, 12, 0.4, 9, False, RGBColor(140, 140, 160))

    num_insights = len(insights)
    risk_count = sum(1 for i in insights if i.get("type") == "risk")
    opp_count = sum(1 for i in insights if i.get("type") == "opportunity")

    summary_items = [
        (f"{num_insights}", "Total Insights", blue),
        (f"{risk_count}", "Risks Identified", red),
        (f"{opp_count}", "Opportunities", green),
        (f"{rows:,}", "Records Analyzed", dark),
    ]
    for idx, (val, label, color) in enumerate(summary_items):
        col_x = 0.4 + idx * 3.2
        _add_rect(slide, col_x, 1.4, 2.9, 1.8, light)
        _add_text_box(slide, val, col_x + 0.1, 1.5, 2.7, 1.0, 32, True, color, PP_ALIGN.CENTER)
        _add_text_box(slide, label, col_x + 0.1, 2.5, 2.7, 0.5, 11, False, dark, PP_ALIGN.CENTER)

    # Summary text
    _add_text_box(slide, (
        f"This AI-powered analysis examined {rows:,} business records to identify performance trends, "
        f"risks, and growth opportunities. {opp_count} actionable opportunities were discovered alongside "
        f"{risk_count} critical risk factors that require attention."
    ), 0.4, 3.5, 12.5, 1.5, 11, False, dark)

    # ── Slide 3: KPI Dashboard ─────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _add_rect(slide, 0, 0, 13.33, 1.1, blue)
    _add_text_box(slide, "KPI Dashboard", 0.4, 0.2, 10, 0.8, 24, True, white)

    numeric_summary = analysis.get("numeric_summary", {})
    kpi_items = list(numeric_summary.items())[:6]
    cols_per_row = 3
    for idx, (col, stats) in enumerate(kpi_items):
        row = idx // cols_per_row
        col_pos = idx % cols_per_row
        x = 0.3 + col_pos * 4.3
        y = 1.3 + row * 2.3
        _add_rect(slide, x, y, 4.0, 2.0, light)
        _add_text_box(slide, col[:18], x + 0.1, y + 0.1, 3.8, 0.5, 10, True, dark)
        _add_text_box(slide, f"${stats.get('total', 0):,.0f}", x + 0.1, y + 0.5, 3.8, 0.7, 18, True, blue, PP_ALIGN.CENTER)
        _add_text_box(slide, f"Avg: {stats.get('mean', 0):,.2f}  |  Max: {stats.get('max', 0):,.2f}", x + 0.1, y + 1.3, 3.8, 0.5, 9, False, dark, PP_ALIGN.CENTER)

    # ── Slide 4: Top Insights ─────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _add_rect(slide, 0, 0, 13.33, 1.1, blue)
    _add_text_box(slide, "Key Business Insights", 0.4, 0.2, 10, 0.8, 24, True, white)

    type_colors = {"risk": red, "opportunity": green, "performance": blue}
    for idx, insight in enumerate(insights[:4], 0):
        y_pos = 1.3 + idx * 1.5
        icolor = type_colors.get(insight.get("type", "info"), blue)
        _add_rect(slide, 0.3, y_pos, 0.07, 1.2, icolor)
        _add_text_box(slide, insight.get("title", "")[:60], 0.5, y_pos, 12, 0.5, 12, True, dark)
        obs = insight.get("observation", "")[:120]
        _add_text_box(slide, obs, 0.5, y_pos + 0.45, 12, 0.8, 9, False, RGBColor(80, 80, 100))

    # ── Slide 5: Recommendations ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _add_rect(slide, 0, 0, 13.33, 1.1, blue)
    _add_text_box(slide, "AI Recommendations", 0.4, 0.2, 10, 0.8, 24, True, white)

    actions = [i.get("action", "") for i in insights if i.get("action")][:5]
    for idx, action in enumerate(actions, 1):
        y_pos = 1.3 + (idx - 1) * 1.1
        _add_rect(slide, 0.3, y_pos + 0.1, 0.5, 0.5, blue)
        _add_text_box(slide, str(idx), 0.3, y_pos + 0.05, 0.5, 0.5, 14, True, white, PP_ALIGN.CENTER)
        _add_text_box(slide, action[:120], 1.0, y_pos, 11.8, 0.8, 11, False, dark)

    # ── Slide 6: Thank You ────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _add_rect(slide, 0, 0, 13.33, 7.5, dark)
    _add_text_box(slide, "Thank You", 1, 2.8, 11, 1.0, 36, True, white, PP_ALIGN.CENTER)
    _add_text_box(slide, f"Generated by {company_name}  |  {generated_at}", 1, 4.0, 11, 0.6, 14, False, RGBColor(160, 175, 220), PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ─────────────────────────────────────────────────────────────────────────────
# EXCEL REPORT
# ─────────────────────────────────────────────────────────────────────────────

def generate_excel_report(
    df: pd.DataFrame,
    analysis: dict,
    insights: list[dict],
    file_id: str,
    company_name: str = "BizInsight AI",
) -> bytes:
    try:
        from openpyxl import Workbook
        from openpyxl.styles import (
            Font, PatternFill, Alignment, Border, Side, numbers
        )
        from openpyxl.utils import get_column_letter
    except ImportError:
        raise ImportError("openpyxl not installed. Run: pip install openpyxl")

    wb = Workbook()
    generated_at = datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC")

    # Color fills
    blue_fill = PatternFill("solid", fgColor="184FB7")
    dark_fill = PatternFill("solid", fgColor="181A2E")
    green_fill = PatternFill("solid", fgColor="0CB282")
    red_fill = PatternFill("solid", fgColor="E73939")
    light_fill = PatternFill("solid", fgColor="EEF0FF")
    alt_fill = PatternFill("solid", fgColor="F5F6FF")
    white_fill = PatternFill("solid", fgColor="FFFFFF")

    header_font = Font(name="Calibri", bold=True, color="FFFFFF", size=11)
    title_font = Font(name="Calibri", bold=True, color="184FB7", size=14)
    body_font = Font(name="Calibri", size=10)
    bold_font = Font(name="Calibri", bold=True, size=10)
    center_align = Alignment(horizontal="center", vertical="center", wrap_text=True)
    left_align = Alignment(horizontal="left", vertical="center", wrap_text=True)

    thin_border = Border(
        left=Side(style="thin", color="CCCCDD"),
        right=Side(style="thin", color="CCCCDD"),
        top=Side(style="thin", color="CCCCDD"),
        bottom=Side(style="thin", color="CCCCDD")
    )

    # ── Sheet 1: Executive Summary ─────────────────────────────────────────────
    ws = wb.active
    ws.title = "Executive Summary"
    ws.column_dimensions["A"].width = 25
    ws.column_dimensions["B"].width = 35
    ws.column_dimensions["C"].width = 20

    ws["A1"] = f"{company_name} — AI Executive Report"
    ws["A1"].font = Font(name="Calibri", bold=True, color="184FB7", size=16)
    ws["A2"] = f"Generated: {generated_at}"
    ws["A2"].font = Font(name="Calibri", color="888899", size=9, italic=True)

    ws["A4"] = "Dataset Overview"
    ws["A4"].font = title_font
    summary_data = [
        ("Industry", analysis.get("industry", "N/A")),
        ("Total Records", f"{analysis.get('shape', {}).get('rows', len(df)):,}"),
        ("Columns", f"{analysis.get('shape', {}).get('columns', len(df.columns))}"),
        ("Analysis Date", generated_at),
        ("Total Insights", str(len(insights))),
        ("Risk Factors", str(sum(1 for i in insights if i.get("type") == "risk"))),
        ("Opportunities", str(sum(1 for i in insights if i.get("type") == "opportunity"))),
    ]
    for row_idx, (label, value) in enumerate(summary_data, start=5):
        ws.cell(row_idx, 1, label).font = bold_font
        ws.cell(row_idx, 2, value).font = body_font
        ws.cell(row_idx, 1).fill = light_fill if row_idx % 2 == 0 else white_fill
        ws.cell(row_idx, 2).fill = light_fill if row_idx % 2 == 0 else white_fill

    # ── Sheet 2: KPI Data ──────────────────────────────────────────────────────
    ws2 = wb.create_sheet("KPI Data")
    ws2.column_dimensions["A"].width = 22
    for col_letter in ["B", "C", "D", "E", "F", "G"]:
        ws2.column_dimensions[col_letter].width = 16

    headers = ["Metric", "Total", "Mean", "Median", "Std Dev", "Min", "Max", "P90"]
    for c_idx, h in enumerate(headers, 1):
        cell = ws2.cell(1, c_idx, h)
        cell.fill = blue_fill
        cell.font = header_font
        cell.alignment = center_align
        cell.border = thin_border

    numeric_summary = analysis.get("numeric_summary", {})
    for r_idx, (col_name, stats) in enumerate(numeric_summary.items(), start=2):
        fill = light_fill if r_idx % 2 == 0 else white_fill
        row_data = [
            col_name,
            stats.get("total", 0),
            stats.get("mean", 0),
            stats.get("median", 0),
            stats.get("std", 0),
            stats.get("min", 0),
            stats.get("max", 0),
            stats.get("p90", 0),
        ]
        for c_idx, val in enumerate(row_data, 1):
            cell = ws2.cell(r_idx, c_idx, val)
            cell.fill = fill
            cell.font = body_font
            cell.border = thin_border
            cell.alignment = center_align if c_idx > 1 else left_align

    # ── Sheet 3: Insights ──────────────────────────────────────────────────────
    ws3 = wb.create_sheet("AI Insights")
    ws3.column_dimensions["A"].width = 10
    ws3.column_dimensions["B"].width = 30
    ws3.column_dimensions["C"].width = 45
    ws3.column_dimensions["D"].width = 45
    ws3.column_dimensions["E"].width = 45

    ins_headers = ["Type", "Title", "Observation", "Interpretation", "Recommended Action"]
    for c_idx, h in enumerate(ins_headers, 1):
        cell = ws3.cell(1, c_idx, h)
        cell.fill = blue_fill
        cell.font = header_font
        cell.alignment = center_align
        cell.border = thin_border

    type_fills = {"risk": red_fill, "opportunity": green_fill, "performance": blue_fill}
    for r_idx, insight in enumerate(insights, start=2):
        itype = insight.get("type", "info")
        fill = type_fills.get(itype, light_fill)
        row_data = [
            itype.upper(),
            insight.get("title", ""),
            insight.get("observation", ""),
            insight.get("interpretation", ""),
            insight.get("action", ""),
        ]
        for c_idx, val in enumerate(row_data, 1):
            cell = ws3.cell(r_idx, c_idx, val)
            cell.fill = fill if c_idx == 1 else (light_fill if r_idx % 2 == 0 else white_fill)
            cell.font = Font(name="Calibri", bold=(c_idx == 1), color=("FFFFFF" if c_idx == 1 else "333344"), size=10)
            cell.border = thin_border
            cell.alignment = Alignment(horizontal="left", vertical="top", wrap_text=True)
        ws3.row_dimensions[r_idx].height = 55

    # ── Sheet 4: Raw Data (first 500 rows) ─────────────────────────────────────
    ws4 = wb.create_sheet("Raw Data")
    sample = df.head(500)
    for c_idx, col_name in enumerate(sample.columns, 1):
        cell = ws4.cell(1, c_idx, col_name)
        cell.fill = dark_fill
        cell.font = header_font
        cell.alignment = center_align
        ws4.column_dimensions[get_column_letter(c_idx)].width = max(12, min(25, len(str(col_name)) + 4))

    for r_idx, row in enumerate(sample.itertuples(index=False), start=2):
        for c_idx, val in enumerate(row, 1):
            cell = ws4.cell(r_idx, c_idx, val)
            cell.font = body_font
            cell.fill = alt_fill if r_idx % 2 == 0 else white_fill
            cell.border = thin_border

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()
